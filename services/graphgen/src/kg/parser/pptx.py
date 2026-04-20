"""PowerPoint parser (.pptx) using python-pptx."""

from __future__ import annotations

from pathlib import Path

from kg.parser import BaseParser, ParsedChunk, ParsedDocument


class PptxParser(BaseParser):
    extensions = ("pptx",)
    mime_types = ("application/vnd.openxmlformats-officedocument.presentationml.presentation",)

    def parse(self, path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation  # type: ignore[import]
        except ImportError as exc:
            raise ImportError("python-pptx is required: pip install python-pptx") from exc

        doc = ParsedDocument.from_path(path, self.mime_types[0])
        prs = Presentation(str(path))
        chunks: list[ParsedChunk] = []

        for slide_num, slide in enumerate(prs.slides):
            parts: list[str] = []
            heading: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # Treat the first text box on a slide as the slide heading
                if not heading:
                    heading = [f"Slide {slide_num + 1}: {text[:80]}"]
                else:
                    parts.append(text)

            combined = " ".join(parts)
            if combined.strip() or heading:
                chunks.append(
                    ParsedChunk.make(doc.doc_id, slide_num, combined or heading[0], heading)
                )

        doc.chunks = chunks
        return doc
